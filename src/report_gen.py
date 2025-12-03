"""
Report generation module for creating PDF and PPTX reports.
"""
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from typing import Dict, List, Any
import os
from datetime import datetime


def create_pdf_report(
    summary: Dict[str, Any],
    insights: Dict[str, List[str]],
    chart_paths: Dict[str, str],
    output_path: str,
    title: str = "AdTech Campaign Performance Report"
) -> str:
    """
    Create a professional PDF report.
    
    Args:
        summary: Dictionary of summary metrics
        insights: Dictionary of AI-generated insights
        chart_paths: Dictionary of chart file paths
        output_path: Path to save the PDF
        title: Report title
        
    Returns:
        Path to saved PDF
    """
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#2E86AB'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#A23B72'),
        spaceAfter=12,
        spaceBefore=12,
        fontName='Helvetica-Bold'
    )
    
    # Title page
    story.append(Spacer(1, 1.5*inch))
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph(f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", 
                          styles['Normal']))
    story.append(PageBreak())
    
    # Executive Summary - KPIs
    story.append(Paragraph("Executive Summary - Key Performance Indicators", heading_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Create KPI table
    kpi_data = [['Metric', 'Value']]
    
    # Format KPIs nicely
    kpi_mapping = {
        'total_impressions': ('Total Impressions', lambda x: f"{x:,.0f}"),
        'total_clicks': ('Total Clicks', lambda x: f"{x:,.0f}"),
        'total_spend': ('Total Spend', lambda x: f"${x:,.2f}"),
        'total_conversions': ('Total Conversions', lambda x: f"{x:,.0f}"),
        'total_revenue': ('Total Revenue', lambda x: f"${x:,.2f}"),
        'overall_CTR': ('CTR', lambda x: f"{x:.2f}%"),
        'overall_CPC': ('CPC', lambda x: f"${x:.2f}"),
        'overall_CPM': ('CPM', lambda x: f"${x:.2f}"),
        'overall_Conversion_Rate': ('Conversion Rate', lambda x: f"{x:.2f}%"),
        'overall_ROAS': ('ROAS', lambda x: f"{x:.2f}x")
    }
    
    for key, (label, formatter) in kpi_mapping.items():
        if key in summary:
            kpi_data.append([label, formatter(summary[key])])
    
    kpi_table = Table(kpi_data, colWidths=[3*inch, 2*inch])
    kpi_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E86AB')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F0F0F0')])
    ]))
    
    story.append(kpi_table)
    story.append(PageBreak())
    
    # AI Insights
    story.append(Paragraph("AI-Powered Insights & Analysis", heading_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Key Insights
    if insights.get('key_insights'):
        story.append(Paragraph("🔍 Key Insights", styles['Heading3']))
        for i, insight in enumerate(insights['key_insights'], 1):
            story.append(Paragraph(f"{i}. {insight}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        story.append(Spacer(1, 0.2*inch))
    
    # Trends
    if insights.get('trends'):
        story.append(Paragraph("📈 Trend Analysis", styles['Heading3']))
        for i, trend in enumerate(insights['trends'], 1):
            story.append(Paragraph(f"{i}. {trend}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        story.append(Spacer(1, 0.2*inch))
    
    # Issues
    if insights.get('issues'):
        story.append(Paragraph("⚠️ Performance Issues", styles['Heading3']))
        for i, issue in enumerate(insights['issues'], 1):
            story.append(Paragraph(f"{i}. {issue}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        story.append(Spacer(1, 0.2*inch))
    
    # Recommendations
    if insights.get('recommendations'):
        story.append(Paragraph("💡 Actionable Recommendations", styles['Heading3']))
        for i, rec in enumerate(insights['recommendations'], 1):
            story.append(Paragraph(f"{i}. {rec}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
    
    story.append(PageBreak())
    
    # Charts
    story.append(Paragraph("Performance Visualizations", heading_style))
    story.append(Spacer(1, 0.2*inch))
    
    for chart_name, chart_path in chart_paths.items():
        if os.path.exists(chart_path):
            story.append(Image(chart_path, width=6*inch, height=3.6*inch))
            story.append(Spacer(1, 0.3*inch))
    
    # Build PDF
    doc.build(story)
    return output_path


def create_pptx_report(
    summary: Dict[str, Any],
    insights: Dict[str, List[str]],
    chart_paths: Dict[str, str],
    output_path: str,
    title: str = "AdTech Campaign Performance Report"
) -> str:
    """
    Create a professional PowerPoint presentation.
    
    Args:
        summary: Dictionary of summary metrics
        insights: Dictionary of AI-generated insights
        chart_paths: Dictionary of chart file paths
        output_path: Path to save the PPTX
        title: Report title
        
    Returns:
        Path to saved PPTX
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title_shape.text = title
    subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    
    # Format title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 134, 171)
    
    # KPI Summary slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Performance Indicators"
    
    tf = body_shape.text_frame
    tf.clear()
    
    # Add KPIs
    kpi_items = []
    if 'overall_CTR' in summary:
        kpi_items.append(f"CTR: {summary['overall_CTR']:.2f}%")
    if 'overall_CPC' in summary:
        kpi_items.append(f"CPC: ${summary['overall_CPC']:.2f}")
    if 'overall_CPM' in summary:
        kpi_items.append(f"CPM: ${summary['overall_CPM']:.2f}")
    if 'overall_Conversion_Rate' in summary:
        kpi_items.append(f"Conversion Rate: {summary['overall_Conversion_Rate']:.2f}%")
    if 'overall_ROAS' in summary:
        kpi_items.append(f"ROAS: {summary['overall_ROAS']:.2f}x")
    if 'total_spend' in summary:
        kpi_items.append(f"Total Spend: ${summary['total_spend']:,.2f}")
    if 'total_revenue' in summary:
        kpi_items.append(f"Total Revenue: ${summary['total_revenue']:,.2f}")
    
    for item in kpi_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    # Key Insights slide
    if insights.get('key_insights'):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = "🔍 Key Insights"
        
        tf = body_shape.text_frame
        tf.clear()
        
        for insight in insights['key_insights'][:5]:
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
            p.font.size = Pt(14)
    
    # Recommendations slide
    if insights.get('recommendations'):
        slide = prs.slides.add_slide(bullet_slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = "💡 Recommendations"
        
        tf = body_shape.text_frame
        tf.clear()
        
        for rec in insights['recommendations'][:5]:
            p = tf.add_paragraph()
            p.text = rec
            p.level = 0
            p.font.size = Pt(14)
    
    # Chart slides
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    
    for chart_name, chart_path in chart_paths.items():
        if os.path.exists(chart_path):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(0.6)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = chart_name.replace('_', ' ').title()
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(46, 134, 171)
            
            # Add chart image
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(9)
            
            slide.shapes.add_picture(chart_path, left, top, width=width)
    
    # Save presentation
    prs.save(output_path)
    return output_path
